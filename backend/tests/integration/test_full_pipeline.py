"""Comprehensive full pipeline integration tests.

Tests the REAL end-to-end flow:
  upload → parse → chunk → embed → store → search → rerank → verify

Validates:
  1. Multi-format document support (MD, TXT, DOCX, XLSX, PPTX)
  2. Multi-tenant isolation (two KBs don't cross-contaminate)
  3. Delete-before-insert deduplication
  4. Metadata propagation through the entire pipeline

Uses:
  - Real: Qdrant (in-memory), MarkItDown parsing, MarkdownNodeParser chunking
  - Fake: Embedding (deterministic char-freq vectors), Sparse (char-freq), Reranker (Jaccard)
"""

from __future__ import annotations

import logging

import pytest
from qdrant_client import AsyncQdrantClient
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker, create_async_engine

from app.models.base import Base
from app.models.document import DocumentStatus
from app.schemas.common import TaskStatus
from app.services.chunking_service import ChunkingService
from app.services.document_service import DocumentService
from app.services.kb_service import KBService
from app.services.parsing_service import ParsingService
from app.services.pipeline_service import PipelineService
from app.services.task_manager import TaskManager
from app.services.vector_store_service import VectorStoreService

logger = logging.getLogger(__name__)


# ─────────────────────────── Fake services (deterministic, no network) ───────


class FakeEmbeddingService:
    """Deterministic dense embedding using character frequency distribution.

    Texts with similar character distributions produce similar vectors,
    so Chinese tech docs and Chinese finance docs will be distinguishable.
    """

    dimension = 1024

    async def embed_texts(self, texts: list[str]) -> list[list[float]]:
        return [self._text_to_vec(t) for t in texts]

    async def embed_query(self, query: str) -> list[float]:
        return self._text_to_vec(query)

    def _text_to_vec(self, text: str) -> list[float]:
        vec = [0.0] * self.dimension
        for ch in text:
            vec[ord(ch) % self.dimension] += 1.0
        norm = sum(x * x for x in vec) ** 0.5 or 1.0
        return [x / norm for x in vec]


class FakeSparseEmbeddingService:
    """Deterministic sparse embedding using character indices."""

    async def embed_texts_async(self, texts: list[str]) -> list[dict]:
        return [self._text_to_sparse(t) for t in texts]

    async def embed_query_async(self, query: str) -> dict:
        return self._text_to_sparse(query)

    def _text_to_sparse(self, text: str) -> dict:
        freq: dict[int, float] = {}
        for ch in text:
            idx = ord(ch) % 5000
            freq[idx] = freq.get(idx, 0) + 1.0
        sorted_idx = sorted(freq.keys())
        return {
            "indices": sorted_idx,
            "values": [freq[i] for i in sorted_idx],
        }


class FakeRerankerService:
    """Reranker based on character-level Jaccard similarity.

    Good enough to rank texts that share more characters with the query higher.
    """

    async def rerank(
        self, query: str, texts: list[str], top_n: int = 3
    ) -> list[dict]:
        if not texts:
            return []
        query_chars = set(query)
        scored = []
        for i, text in enumerate(texts):
            text_chars = set(text)
            union = len(query_chars | text_chars) or 1
            score = len(query_chars & text_chars) / union
            scored.append({"index": i, "score": score, "text": text})
        scored.sort(key=lambda x: x["score"], reverse=True)
        return scored[:top_n]


# ─────────────────────────── Test file generators ────────────────────────────


def create_md_file(path, content: str) -> None:
    path.write_text(content, encoding="utf-8")


def create_txt_file(path, content: str) -> None:
    path.write_text(content, encoding="utf-8")


def create_docx_file(path, paragraphs: list[str]) -> None:
    from docx import Document as DocxDocument

    doc = DocxDocument()
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save(str(path))


def create_xlsx_file(path, rows: list[list[str]]) -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    wb.save(str(path))


def create_pptx_file(path, slides: list[str]) -> None:
    from pptx import Presentation

    prs = Presentation()
    for title in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = title
    prs.save(str(path))


# ─────────────────────────── Fixtures ────────────────────────────────────────


@pytest.fixture
async def db_engine():
    engine = create_async_engine("sqlite+aiosqlite://", echo=False)
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
    yield engine
    await engine.dispose()


@pytest.fixture
async def session_factory(db_engine):
    return async_sessionmaker(db_engine, class_=AsyncSession, expire_on_commit=False)


@pytest.fixture
async def qdrant():
    client = AsyncQdrantClient(location=":memory:")
    yield client
    await client.close()


@pytest.fixture
def services(qdrant):
    """Bundle of all services needed for the pipeline."""
    return {
        "parsing": ParsingService(),
        "chunking": ChunkingService(chunk_size=256, chunk_overlap=32),
        "embedding": FakeEmbeddingService(),
        "sparse": FakeSparseEmbeddingService(),
        "vector_store": VectorStoreService(qdrant),
        "reranker": FakeRerankerService(),
        "task_manager": TaskManager(),
    }


async def run_pipeline_sync(
    session: AsyncSession,
    services: dict,
    file_path: str,
    doc_id: str,
    file_name: str,
    kb_id: str,
) -> str:
    """Helper: run the full pipeline and return the task_id."""
    tm: TaskManager = services["task_manager"]
    task_info = tm.create_task()

    pipeline = PipelineService(
        session=session,
        parsing_service=services["parsing"],
        chunking_service=services["chunking"],
        embedding_service=services["embedding"],
        sparse_embedding_service=services["sparse"],
        vector_store_service=services["vector_store"],
        task_manager=tm,
    )

    await pipeline.run_pipeline(
        task_id=task_info.task_id,
        file_path=file_path,
        doc_id=doc_id,
        file_name=file_name,
        knowledge_base_id=kb_id,
    )

    return task_info.task_id


# ─────────────────────────── Tests ───────────────────────────────────────────


class TestMultiFormatPipeline:
    """Verify that MD, TXT, DOCX, XLSX, PPTX all go through the full pipeline."""

    async def test_md_pipeline(self, session_factory, services, tmp_path):
        kb_id = "kb_md_test"
        kb_id = "kb_md_test"
        await services["vector_store"].create_collection(kb_id)

        md_file = tmp_path / "router_manual.md"
        create_md_file(
            md_file,
            "# 路由器手册\n\n## XJ-998 规格\n\nXJ-998路由器支持Wi-Fi 7技术。\n\n"
            "## 故障排除\n\n如果路由器LED灯闪红灯，请按住重置键10秒。\n",
        )

        async with session_factory() as session:
            doc_svc = DocumentService(session)
            # Create a KB record that the doc can reference
            kb_svc = KBService(session)
            try:
                await kb_svc.create("MD Test KB")
            except Exception:
                pass
            await doc_svc.create("doc_md_001", "router_manual.md", kb_id)
            task_id = await run_pipeline_sync(
                session, services, str(md_file), "doc_md_001", "router_manual.md", kb_id
            )

        tm = services["task_manager"]
        task = tm.get_task(task_id)
        assert task.status == TaskStatus.COMPLETED, f"Pipeline failed: {task.error}"
        assert task.result["chunk_count"] > 0

        # Verify doc status in DB
        async with session_factory() as session:
            doc_svc = DocumentService(session)
            doc = await doc_svc.get_by_doc_id_and_kb("doc_md_001", kb_id)
            assert doc.status == DocumentStatus.COMPLETED
            assert doc.chunk_count > 0

        # Verify vectors in Qdrant — search should return results
        query_vec = await services["embedding"].embed_query("XJ-998路由器")
        query_sparse = await services["sparse"].embed_query_async("XJ-998路由器")
        results = await services["vector_store"].hybrid_search(
            kb_id, query_vec, query_sparse, top_k=5
        )
        assert len(results) > 0
        texts = [r["payload"]["text"] for r in results]
        assert any("XJ-998" in t or "路由器" in t for t in texts)
        # Verify metadata
        for r in results:
            assert r["payload"]["doc_id"] == "doc_md_001"
            assert r["payload"]["file_name"] == "router_manual.md"
            assert r["payload"]["knowledge_base_id"] == kb_id

    async def test_txt_pipeline(self, session_factory, services, tmp_path):
        kb_id = "kb_txt_test"
        await services["vector_store"].create_collection(kb_id)

        txt_file = tmp_path / "expense_policy.txt"
        create_txt_file(
            txt_file,
            "公司报销制度\n\n员工差旅报销必须在出差结束后14天内提交。"
            "超过30天未提交的报销单将自动作废。\n"
            "报销金额超过5000元需部门经理审批。\n",
        )

        async with session_factory() as session:
            kb_svc = KBService(session)
            await kb_svc.create("TXT Test KB")
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_txt_001", "expense_policy.txt", kb_id)
            task_id = await run_pipeline_sync(
                session, services, str(txt_file), "doc_txt_001", "expense_policy.txt", kb_id
            )

        task = services["task_manager"].get_task(task_id)
        assert task.status == TaskStatus.COMPLETED, f"Pipeline failed: {task.error}"
        assert task.result["chunk_count"] > 0

    async def test_docx_pipeline(self, session_factory, services, tmp_path):
        kb_id = "kb_docx_test"
        await services["vector_store"].create_collection(kb_id)

        docx_file = tmp_path / "network_guide.docx"
        create_docx_file(
            docx_file,
            [
                "网络配置指南",
                "第一步：连接网线到路由器WAN口。",
                "第二步：打开浏览器访问192.168.1.1进行初始配置。",
                "第三步：设置WiFi名称和密码。",
            ],
        )

        async with session_factory() as session:
            kb_svc = KBService(session)
            await kb_svc.create("DOCX Test KB")
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_docx_001", "network_guide.docx", kb_id)
            task_id = await run_pipeline_sync(
                session, services, str(docx_file), "doc_docx_001", "network_guide.docx", kb_id
            )

        task = services["task_manager"].get_task(task_id)
        assert task.status == TaskStatus.COMPLETED, f"Pipeline failed: {task.error}"
        assert task.result["chunk_count"] > 0

        # Verify content is searchable
        query_vec = await services["embedding"].embed_query("路由器网络配置")
        query_sparse = await services["sparse"].embed_query_async("路由器网络配置")
        results = await services["vector_store"].hybrid_search(
            kb_id, query_vec, query_sparse, top_k=5
        )
        assert len(results) > 0

    async def test_xlsx_pipeline(self, session_factory, services, tmp_path):
        kb_id = "kb_xlsx_test"
        await services["vector_store"].create_collection(kb_id)

        xlsx_file = tmp_path / "budget.xlsx"
        create_xlsx_file(
            xlsx_file,
            [
                ["部门", "预算(万)", "实际支出(万)", "结余(万)"],
                ["研发部", "500", "420", "80"],
                ["市场部", "300", "310", "-10"],
                ["人力资源", "200", "180", "20"],
            ],
        )

        async with session_factory() as session:
            kb_svc = KBService(session)
            await kb_svc.create("XLSX Test KB")
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_xlsx_001", "budget.xlsx", kb_id)
            task_id = await run_pipeline_sync(
                session, services, str(xlsx_file), "doc_xlsx_001", "budget.xlsx", kb_id
            )

        task = services["task_manager"].get_task(task_id)
        assert task.status == TaskStatus.COMPLETED, f"Pipeline failed: {task.error}"
        assert task.result["chunk_count"] > 0

    async def test_pptx_pipeline(self, session_factory, services, tmp_path):
        kb_id = "kb_pptx_test"
        await services["vector_store"].create_collection(kb_id)

        pptx_file = tmp_path / "product_intro.pptx"
        create_pptx_file(
            pptx_file,
            ["产品发布会", "XJ-998新品路由器", "核心特性：Wi-Fi 7"],
        )

        async with session_factory() as session:
            kb_svc = KBService(session)
            await kb_svc.create("PPTX Test KB")
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_pptx_001", "product_intro.pptx", kb_id)
            task_id = await run_pipeline_sync(
                session, services, str(pptx_file), "doc_pptx_001", "product_intro.pptx", kb_id
            )

        task = services["task_manager"].get_task(task_id)
        assert task.status == TaskStatus.COMPLETED, f"Pipeline failed: {task.error}"
        assert task.result["chunk_count"] > 0


class TestMultiTenantIsolation:
    """Verify that different knowledge bases are completely isolated.

    Creates two KBs with domain-specific content, then queries each
    to confirm zero cross-contamination.
    """

    async def test_two_kbs_no_crosstalk(self, session_factory, services, tmp_path):
        # ── Setup: two completely different knowledge domains ──

        kb_tech = "kb_tech_support"
        kb_finance = "kb_finance_dept"

        await services["vector_store"].create_collection(kb_tech)
        await services["vector_store"].create_collection(kb_finance)

        # Tech KB documents
        tech_md = tmp_path / "router_specs.md"
        create_md_file(
            tech_md,
            "# XJ-998路由器技术规格\n\n"
            "## 无线标准\n\nXJ-998支持最新的Wi-Fi 7 (802.11be)标准，"
            "最高速率可达36Gbps。采用MLO多链路聚合技术。\n\n"
            "## 硬件参数\n\n- CPU: 四核2.0GHz\n- 内存: 1GB DDR4\n"
            "- 天线: 8根高增益天线\n",
        )
        tech_txt = tmp_path / "troubleshooting.txt"
        create_txt_file(
            tech_txt,
            "路由器常见故障排查指南\n\n"
            "问题一：红灯闪烁\n解决方案：按住reset键10秒恢复出厂设置。\n\n"
            "问题二：网速慢\n解决方案：检查信道是否拥挤，切换至5GHz频段。\n",
        )

        # Finance KB documents
        finance_md = tmp_path / "expense_rules.md"
        create_md_file(
            finance_md,
            "# 公司财务报销制度\n\n"
            "## 差旅报销\n\n员工出差需提前申请审批。住宿标准：一线城市500元/晚，"
            "二线城市350元/晚。\n\n"
            "## 餐饮补贴\n\n午餐补贴标准为每人每日30元，加班餐另报。\n",
        )
        finance_xlsx = tmp_path / "q4_budget.xlsx"
        create_xlsx_file(
            finance_xlsx,
            [
                ["科目", "Q4预算", "Q4实际", "差异"],
                ["办公用品", "50000", "48000", "2000"],
                ["差旅费", "200000", "185000", "15000"],
                ["培训费", "100000", "95000", "5000"],
            ],
        )

        # ── Ingest: tech docs → kb_tech, finance docs → kb_finance ──

        async with session_factory() as session:
            kb_svc = KBService(session)
            await kb_svc.create("Tech Support", "技术支持知识库")
            await kb_svc.create("Finance Dept", "财务部知识库")

        # Ingest tech docs
        async with session_factory() as session:
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_tech_md", "router_specs.md", kb_tech)
            await run_pipeline_sync(
                session, services, str(tech_md), "doc_tech_md", "router_specs.md", kb_tech
            )
        async with session_factory() as session:
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_tech_txt", "troubleshooting.txt", kb_tech)
            await run_pipeline_sync(
                session, services, str(tech_txt), "doc_tech_txt", "troubleshooting.txt", kb_tech
            )

        # Ingest finance docs
        async with session_factory() as session:
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_fin_md", "expense_rules.md", kb_finance)
            await run_pipeline_sync(
                session, services, str(finance_md), "doc_fin_md", "expense_rules.md", kb_finance
            )
        async with session_factory() as session:
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_fin_xlsx", "q4_budget.xlsx", kb_finance)
            await run_pipeline_sync(
                session, services, str(finance_xlsx), "doc_fin_xlsx", "q4_budget.xlsx", kb_finance
            )

        # ── Query: verify isolation ──

        embed = services["embedding"]
        sparse = services["sparse"]
        vs = services["vector_store"]
        reranker = services["reranker"]

        # --- Test 1: Tech query on tech KB → should find router content ---
        q1 = "XJ-998路由器支持什么无线标准"
        q1_dense = await embed.embed_query(q1)
        q1_sparse = await sparse.embed_query_async(q1)
        tech_results = await vs.hybrid_search(kb_tech, q1_dense, q1_sparse, top_k=10)
        assert len(tech_results) > 0, "Tech KB should return results for tech query"

        tech_texts = [r["payload"]["text"] for r in tech_results]
        reranked_tech = await reranker.rerank(q1, tech_texts, top_n=3)
        assert len(reranked_tech) > 0

        # ALL results must belong to tech KB
        for r in tech_results:
            assert r["payload"]["knowledge_base_id"] == kb_tech, (
                f"ISOLATION VIOLATION: tech KB returned doc from {r['payload']['knowledge_base_id']}"
            )
            assert r["payload"]["doc_id"].startswith("doc_tech"), (
                f"ISOLATION VIOLATION: tech KB returned doc_id={r['payload']['doc_id']}"
            )

        # --- Test 2: Finance query on finance KB → should find finance content ---
        q2 = "差旅报销标准住宿费"
        q2_dense = await embed.embed_query(q2)
        q2_sparse = await sparse.embed_query_async(q2)
        fin_results = await vs.hybrid_search(kb_finance, q2_dense, q2_sparse, top_k=10)
        assert len(fin_results) > 0, "Finance KB should return results for finance query"

        # ALL results must belong to finance KB
        for r in fin_results:
            assert r["payload"]["knowledge_base_id"] == kb_finance, (
                f"ISOLATION VIOLATION: finance KB returned doc from "
                f"{r['payload']['knowledge_base_id']}"
            )
            assert r["payload"]["doc_id"].startswith("doc_fin"), (
                f"ISOLATION VIOLATION: finance KB returned doc_id={r['payload']['doc_id']}"
            )

        # --- Test 3: Cross-query – tech query on FINANCE KB → should NOT return tech docs ---
        cross_results = await vs.hybrid_search(kb_finance, q1_dense, q1_sparse, top_k=10)
        for r in cross_results:
            assert r["payload"]["knowledge_base_id"] == kb_finance, (
                "CROSS-CONTAMINATION: finance KB search returned a tech doc!"
            )
            assert "XJ-998" not in r["payload"].get("file_name", ""), (
                "CROSS-CONTAMINATION: finance KB returned router specs file!"
            )

        # --- Test 4: Cross-query – finance query on TECH KB → should NOT return finance docs ---
        cross_results2 = await vs.hybrid_search(kb_tech, q2_dense, q2_sparse, top_k=10)
        for r in cross_results2:
            assert r["payload"]["knowledge_base_id"] == kb_tech, (
                "CROSS-CONTAMINATION: tech KB search returned a finance doc!"
            )
            assert r["payload"]["doc_id"].startswith("doc_tech"), (
                "CROSS-CONTAMINATION: tech KB returned finance doc_id!"
            )

        logger.info("Multi-tenant isolation test PASSED — zero cross-contamination")

    async def test_three_kbs_full_isolation(self, session_factory, services, tmp_path):
        """Test with 3 KBs to ensure N-way isolation, not just 2-way."""
        kb_a = "kb_dept_a"
        kb_b = "kb_dept_b"
        kb_c = "kb_dept_c"

        for kb_id in [kb_a, kb_b, kb_c]:
            await services["vector_store"].create_collection(kb_id)

        docs = {
            kb_a: ("alpha.md", "# Alpha部门\n\n这是Alpha部门的专属内部文件，包含项目管理规范。\n"),
            kb_b: ("beta.md", "# Beta部门\n\n这是Beta部门的测试用例模板和质量标准。\n"),
            kb_c: ("gamma.md", "# Gamma部门\n\n这是Gamma部门的客户服务流程和响应时间要求。\n"),
        }

        async with session_factory() as session:
            kb_svc = KBService(session)
            await kb_svc.create("Dept A")
            await kb_svc.create("Dept B")
            await kb_svc.create("Dept C")

        for kb_id, (fname, content) in docs.items():
            fpath = tmp_path / fname
            create_md_file(fpath, content)
            doc_id = f"doc_{fname.split('.')[0]}"
            async with session_factory() as session:
                doc_svc = DocumentService(session)
                await doc_svc.create(doc_id, fname, kb_id)
                await run_pipeline_sync(session, services, str(fpath), doc_id, fname, kb_id)

        # Each KB should only contain its own data
        embed = services["embedding"]
        sparse = services["sparse"]
        vs = services["vector_store"]

        for kb_id, (fname, _) in docs.items():
            q_dense = await embed.embed_query(fname)
            q_sparse = await sparse.embed_query_async(fname)
            results = await vs.hybrid_search(kb_id, q_dense, q_sparse, top_k=10)
            for r in results:
                assert r["payload"]["knowledge_base_id"] == kb_id, (
                    f"ISOLATION VIOLATION: {kb_id} returned data from "
                    f"{r['payload']['knowledge_base_id']}"
                )


class TestDeleteBeforeInsert:
    """Verify that re-uploading the same doc replaces old vectors completely."""

    async def test_upsert_replaces_old_data(self, session_factory, services, tmp_path):
        kb_id = "kb_upsert_test"
        await services["vector_store"].create_collection(kb_id)

        async with session_factory() as session:
            kb_svc = KBService(session)
            await kb_svc.create("Upsert Test KB")

        # Version 1: original document
        v1_file = tmp_path / "product_v1.md"
        create_md_file(v1_file, "# 产品说明 V1\n\n旧版本的产品描述，包含过时的技术参数。\n")

        async with session_factory() as session:
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_product", "product.md", kb_id)
            await run_pipeline_sync(
                session, services, str(v1_file), "doc_product", "product.md", kb_id
            )

        # Search should find V1 content
        q_dense = await services["embedding"].embed_query("产品说明")
        q_sparse = await services["sparse"].embed_query_async("产品说明")
        v1_results = await services["vector_store"].hybrid_search(
            kb_id, q_dense, q_sparse, top_k=10
        )
        v1_texts = " ".join(r["payload"]["text"] for r in v1_results)
        assert "V1" in v1_texts or "旧版本" in v1_texts

        # Version 2: updated document (same doc_id!)
        v2_file = tmp_path / "product_v2.md"
        create_md_file(v2_file, "# 产品说明 V2\n\n全新版本，支持Wi-Fi 7和最新安全协议。\n")

        async with session_factory() as session:
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_product", "product.md", kb_id)
            await run_pipeline_sync(
                session, services, str(v2_file), "doc_product", "product.md", kb_id
            )

        # Search should now find ONLY V2 content — V1 must be gone
        v2_results = await services["vector_store"].hybrid_search(
            kb_id, q_dense, q_sparse, top_k=10
        )
        for r in v2_results:
            if r["payload"]["doc_id"] == "doc_product":
                assert "旧版本" not in r["payload"]["text"], (
                    "STALE DATA: V1 content still present after V2 upsert!"
                )

        # Verify doc record reflects V2
        async with session_factory() as session:
            doc_svc = DocumentService(session)
            doc = await doc_svc.get_by_doc_id_and_kb("doc_product", kb_id)
            assert doc.status == DocumentStatus.COMPLETED


class TestMetadataPropagation:
    """Verify that metadata survives the entire pipeline and appears in search results."""

    async def test_metadata_in_search_results(self, session_factory, services, tmp_path):
        kb_id = "kb_meta_test"
        await services["vector_store"].create_collection(kb_id)

        async with session_factory() as session:
            kb_svc = KBService(session)
            await kb_svc.create("Metadata Test KB")

        md_file = tmp_path / "meta_test.md"
        create_md_file(
            md_file,
            "# 测试文档\n\n## 第一节\n\n这是用来测试元数据传播的内容。\n\n"
            "## 第二节\n\n另一个章节的内容，用来验证chunk_index递增。\n",
        )

        async with session_factory() as session:
            doc_svc = DocumentService(session)
            await doc_svc.create("doc_meta_001", "meta_test.md", kb_id)
            await run_pipeline_sync(
                session, services, str(md_file), "doc_meta_001", "meta_test.md", kb_id
            )

        q_dense = await services["embedding"].embed_query("测试元数据")
        q_sparse = await services["sparse"].embed_query_async("测试元数据")
        results = await services["vector_store"].hybrid_search(
            kb_id, q_dense, q_sparse, top_k=10
        )

        assert len(results) > 0
        for r in results:
            payload = r["payload"]
            # Required metadata fields
            assert payload["doc_id"] == "doc_meta_001"
            assert payload["file_name"] == "meta_test.md"
            assert payload["knowledge_base_id"] == kb_id
            assert "chunk_index" in payload
            assert isinstance(payload["chunk_index"], int)
            assert "text" in payload
            assert len(payload["text"]) > 0

        # Verify chunk_index values are assigned
        indices = sorted(r["payload"]["chunk_index"] for r in results)
        assert indices == list(range(len(indices))), (
            f"chunk_index should be 0-based sequential, got {indices}"
        )


class TestRetrievalServiceEndToEnd:
    """Test the full retrieval orchestration: embed → search → rerank."""

    async def test_retrieve_returns_ranked_results(self, session_factory, services, tmp_path):
        from app.services.retrieval_service import RetrievalService

        kb_id = "kb_retrieval_e2e"
        await services["vector_store"].create_collection(kb_id)

        async with session_factory() as session:
            kb_svc = KBService(session)
            await kb_svc.create("Retrieval E2E KB")

        # Ingest multiple documents
        docs = [
            ("spec.md", "doc_spec", "# XJ-998规格\n\n支持Wi-Fi 7技术，最大速率36Gbps。\n"),
            ("faq.md", "doc_faq", "# 常见问题\n\n问：如何重置路由器？\n答：按住reset键10秒。\n"),
            (
                "unrelated.md",
                "doc_unrelated",
                "# 公司年会\n\n今年年会定于12月在三亚举办，请大家准备节目。\n",
            ),
        ]
        for fname, doc_id, content in docs:
            fpath = tmp_path / fname
            create_md_file(fpath, content)
            async with session_factory() as session:
                doc_svc = DocumentService(session)
                await doc_svc.create(doc_id, fname, kb_id)
                await run_pipeline_sync(session, services, str(fpath), doc_id, fname, kb_id)

        # Build the retrieval service
        retrieval = RetrievalService(
            embedding_service=services["embedding"],
            sparse_embedding_service=services["sparse"],
            vector_store_service=services["vector_store"],
            reranker_service=services["reranker"],
        )

        result = await retrieval.retrieve(
            knowledge_base_id=kb_id,
            query="XJ-998路由器Wi-Fi 7速率",
            top_k=10,
            top_n=2,
        )

        assert result["total_candidates"] > 0
        assert len(result["source_nodes"]) <= 2
        assert result["top_k_used"] == 10
        assert result["top_n_used"] == 2

        # Each source node must have required fields
        for node in result["source_nodes"]:
            assert "text" in node
            assert "score" in node
            assert "doc_id" in node
            assert "file_name" in node
            assert "knowledge_base_id" in node
            assert node["knowledge_base_id"] == kb_id
